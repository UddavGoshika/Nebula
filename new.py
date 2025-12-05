
import os, io, re, json, shutil
from pathlib import Path
from typing import List
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from sentence_transformers import CrossEncoder
from starlette.requests import Request
from fastapi.responses import RedirectResponse, JSONResponse
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
from rank_bm25 import BM25Okapi
from PyPDF2 import PdfReader
import docx
import google.generativeai as genai
import uvicorn
from authlib.integrations.starlette_client import OAuth
from datetime import datetime
from redis.asyncio import Redis
from pydantic import BaseModel
from dotenv import load_dotenv
import redis
import uuid
import jwt  # already there
import redis
from functools import lru_cache

load_dotenv()  # must be first before reading GEMINI_API_KEY

RERANKER = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2")


UPSTASH_URL = os.getenv("UPSTASH_REDIS_URL")  # ← change to this

@lru_cache()
def get_redis():
    return redis.from_url(UPSTASH_URL, decode_responses=True, socket_keepalive=True, retry_on_timeout=True)  # ← add keepalive + retry

r = get_redis() if UPSTASH_URL else None
# ----------------------------------------------------------
# ✅ CONFIG
# ----------------------------------------------------------
BASE = Path(__file__).parent
DATA = BASE / "data"
UPLOADS = BASE / "uploads"

DATA.mkdir(exist_ok=True)
UPLOADS.mkdir(exist_ok=True)

# Embeddings
EMB = SentenceTransformer("all-MiniLM-L6-v2")
DIM = EMB.get_sentence_embedding_dimension()

# ✅ Correct Gemini model
GEMINI_MODEL = "gemini-2.5-f"

# ✅ Configure Gemini
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    print("❌ ERROR: GEMINI_API_KEY missing — set it before running.")
else:
    genai.configure(api_key=GEMINI_API_KEY)


# ----------------------------------------------------------
# ✅ Storage: FAISS, Metadata, BM25
# ----------------------------------------------------------
INDEX_PATH = DATA / "faiss.index"
META_PATH = DATA / "meta.json"
BM25_PATH = DATA / "bm25.json"

# FAISS
if INDEX_PATH.exists():
    index = faiss.read_index(str(INDEX_PATH))
else:
    index = faiss.IndexFlatL2(DIM)

# Metadata
if META_PATH.exists():
    METADATA = json.loads(META_PATH.read_text(encoding="utf-8"))
else:
    METADATA = []

# BM25
if BM25_PATH.exists():
    token_corpus = json.loads(BM25_PATH.read_text(encoding="utf-8"))
else:
    token_corpus = [re.findall(r"[a-z0-9]+", m["text"].lower()) for m in METADATA]

BM25 = BM25Okapi(token_corpus) if token_corpus else None


def persist():
    faiss.write_index(index, str(INDEX_PATH))
    META_PATH.write_text(json.dumps(METADATA, ensure_ascii=False), encoding="utf-8")
    BM25_PATH.write_text(json.dumps(token_corpus), encoding="utf-8")

# new one 

def get_user_id(request: Request):
    auth = request.headers.get("Authorization")
    if not auth or not auth.startswith("Bearer "):
        return None
    try:
        token = auth.split(" ")[1]
        payload = jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        return payload["sub"]  # this is the Google email
    except:
        return None




# ----------------------------------------------------------
# ✅ Document extraction
# ----------------------------------------------------------
def extract_pdf(raw, name):
    reader = PdfReader(io.BytesIO(raw))
    pages = []
    for i, p in enumerate(reader.pages, start=1):
        t = p.extract_text() or ""
        if t.strip():
            pages.append({"text": t, "source": name, "page": i})
    return pages

import io
import cv2
import numpy as np
from PIL import Image
import pytesseract
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup
import docx   # for extract_docx


# ---------- IMAGE / SCREENSHOT (OCR) ----------

def ocr_terminal_screenshot(raw_bytes):
    img = Image.open(io.BytesIO(raw_bytes)).convert("L")
    img = np.array(img)

    # Upscale a bit so small text is easier to read
    img = cv2.resize(img, None, fx=2.0, fy=2.0, interpolation=cv2.INTER_LINEAR)

    # Binarize (black/white) for better contrast
    img = cv2.threshold(img, 0, 255,
                        cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]

    text = pytesseract.image_to_string(
        img,
        lang="eng",
        config="--oem 3 --psm 6"
    )
    return text


def extract_image(raw, name):
    t = ocr_terminal_screenshot(raw)
    return [{"text": t, "source": name, "page": None}] if t.strip() else []


# ---------- CSV ----------

def extract_csv(raw, name, max_rows=2000):
    df = pd.read_csv(io.BytesIO(raw), nrows=max_rows)
    if df.empty:
        return []
    t = df.to_csv(index=False)
    return [{"text": t, "source": name, "page": None}] if t.strip() else []


# ---------- EXCEL ----------

def extract_excel(raw, name, max_rows=2000):
    xls = pd.ExcelFile(io.BytesIO(raw))
    blocks = []

    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name, nrows=max_rows)
        if df.empty:
            continue
        t = f"Sheet: {sheet_name}\n" + df.to_csv(index=False)
        if t.strip():
            blocks.append({
                "text": t,
                "source": f"{name}#{sheet_name}",
                "page": None
            })

    return blocks


# ---------- POWERPOINT (PPTX) ----------

def extract_pptx(raw, name):
    prs = Presentation(io.BytesIO(raw))
    blocks = []

    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        t = "\n".join(parts)
        if t.strip():
            blocks.append({
                "text": t,
                "source": f"{name}#slide-{i}",
                "page": i
            })

    return blocks


# ---------- HTML / WEB PAGE ----------

def extract_html(raw, name):
    html = raw.decode("utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")

    # Remove scripts/styles
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    title = (soup.title.string or "").strip() if soup.title else ""
    body_text = soup.get_text(separator="\n")
    full_text = (title + "\n\n" + body_text).strip() if title else body_text.strip()

    return [{"text": full_text, "source": name, "page": None}] if full_text else []





def extract_docx(raw, name):
    d = docx.Document(io.BytesIO(raw))
    t = "\n".join(p.text for p in d.paragraphs)
    return [{"text": t, "source": name, "page": None}] if t.strip() else []


def extract_txt(raw, name):
    t = raw.decode("utf-8", errors="ignore")
    return [{"text": t, "source": name, "page": None}] if t.strip() else []




def extract_blocks(f: UploadFile):
    name = f.filename.lower()
    raw = f.file.read()

    if name.endswith(".pdf"): return extract_pdf(raw, f.filename)
    if name.endswith(".jpeg") or name.endswith(".jpg") or name.endswith(".png"):
        return extract_image(raw, f.filename)
    if name.endswith(".csv"): return extract_csv(raw, f.filename)
    if name.endswith(".xlsx") or name.endswith(".xls"): return extract_excel(raw, f.filename)
    if name.endswith(".pptx"): return extract_pptx(raw, f.filename)
    if name.endswith(".html") or name.endswith(".htm"): return extract_html(raw, f.filename)
    if name.endswith(".docx"): return extract_docx(raw, f.filename)
    if name.endswith(".txt"): return extract_txt(raw, f.filename)
    return extract_txt(raw, f.filename)



# ==================== SMART CHUNKING (Replace your current chunk_text function with this)
def smart_chunk(text: str, max_chars: int = 500, overlap: int = 100):
    """Optimized for Groq Llama-3.1-8B → small, clean, precise chunks"""
    if not text or len(text.strip()) < 50:
        return []
    
    chunks = []
    current = ""
    paragraphs = [p.strip() for p in re.split(r'\n\s*\n', text) if p.strip()]

    for para in paragraphs:
        # Split on sentence boundaries but keep proper names/titles intact
        sentences = re.split(r'(?<=[.!?])\s+(?=[A-Z])|(?<=[.!?])\s+', para)
        sentences = [s.strip() for s in sentences if s.strip() and len(s) > 15]

        for sentence in sentences:
            # If adding this sentence overflows → save current chunk
            if len(current) + len(sentence) + 1 > max_chars and current:
                chunks.append(current.strip())
                # Overlap: keep last N chars for context continuity
                current = current[-overlap:] + " " + sentence
            else:
                current = (current + " " + sentence).strip()
        
        # Add paragraph break for natural flow
        if current:
            current += "\n"

    if current.strip():
        chunks.append(current.strip())

    return chunks






# ==================== INGESTION (Replace your ingest function with this)
def ingest(blocks):
    global METADATA, token_corpus, BM25
    new_chunks = []

    for b in blocks:
        for ch in smart_chunk(b["text"]):
            new_chunks.append({"text": ch, "source": b["source"], "page": b["page"]})

    if not new_chunks:
        return 0

    vecs = EMB.encode([c["text"] for c in new_chunks], convert_to_numpy=True)
    index.add(vecs)
    METADATA.extend(new_chunks)

    toks = [re.findall(r"[a-z0-9]+", c["text"].lower()) for c in new_chunks]
    token_corpus.extend(toks)
    BM25 = BM25Okapi(token_corpus)

    persist()
    return len(new_chunks)









# ----------------------------------------------------------
# ✅ Hybrid Retrieval (FAISS + BM25 + rerank)
# ----------------------------------------------------------
# def rerank(query, topk=6):
#     q_emb = EMB.encode([query], convert_to_numpy=True)

#     faiss_idx = index.search(q_emb, 10)[1][0]
#     bm25_idx = np.argsort(BM25.get_scores(re.findall(r"[a-z0-9]+", query.lower())))[::-1][:10]

#     candidates = set([int(i) for i in faiss_idx if i >= 0]) | set(bm25_idx)
#     candidates = list(candidates)

#     if not candidates:
#         return []

#     texts = [METADATA[i]["text"] for i in candidates]
#     embs = EMB.encode(texts, convert_to_numpy=True)

#     q = q_emb[0] / (np.linalg.norm(q_emb[0]) + 1e-9)
#     e = embs / (np.linalg.norm(embs, axis=1, keepdims=True) + 1e-9)

#     sims = e @ q
#     order = np.argsort(-sims)[:topk]

#     return [candidates[i] for i in order]


def rerank(query, topk=3):
    # Step 1: vector candidates from FAISS + BM25 (same as before)
    q_emb = EMB.encode([query], convert_to_numpy=True)
    faiss_idx = index.search(q_emb, 10)[1][0]
    if BM25:
        bm25_idx = np.argsort(BM25.get_scores(re.findall(r"[a-z0-9]+", query.lower())))[::-1][:10]
    else:
        bm25_idx = []

    # bm25_idx = np.argsort(BM25.get_scores(re.findall(r"[a-z0-9]+", query.lower())))[::-1][:10]

    candidates = set([int(i) for i in faiss_idx if i >= 0]) | set(bm25_idx)
    candidates = list(candidates)
    if not candidates:
        return []

    # Step 2: build (query, text) pairs
    pairs = [(query, METADATA[i]["text"]) for i in candidates]

    # Step 3: score with cross-encoder
    scores = RERANKER.predict(pairs)

    # Step 4: sort by score
    order = np.argsort(-scores)[:topk]

    # Step 5: return best chunk ids
    return [candidates[i] for i in order]




# ----------------------------------------------------------
# ✅ Gemini LLM Answer
# ----------------------------------------------------------
# def answer_query(question: str):

#     idxs = rerank(question , topk=3)
#     context = "\n---\n".join(METADATA[i]["text"] for i in idxs)

#     cites = [
#         f"[{n+1}] {METADATA[i]['source']} p.{METADATA[i]['page'] or '-'}"
#         for n, i in enumerate(idxs)
#     ]
    
#     MAX_CONTEXT = 5000
#     if len(context) > MAX_CONTEXT:
#         context = context[:MAX_CONTEXT] + "\n\n[Context truncated]"


#     prompt = f"""
# Use the context to answer.

# CONTEXT:
# {context}

# QUESTION:
# {question}

# ANSWER:
# """
#     print("📝 Prompt to Gemini:\n", prompt)  # <--- DEBUG PROMPT
#     model = genai.GenerativeModel(GEMINI_MODEL)
#     response = model.generate_content(prompt)
#     print("✅ Gemini called!")   # <--- PROOF IN TERMINAL
#     # print("📝 Gemini Response:\n", response.text)  # <--- DEBUG RESPONSE
#     answer = response.text.strip()
#     # if "Sources:" not in answer:
#     #     answer += "\n\nSources:\n" + "\n".join(cites)

#     return answer



# def answer_hybrid_llm(question: str, role: str = "User"):
#     """Hybrid mode: FAISS + BM25 + RRF + Reranker + LLM synthesis"""
#     q_emb = EMB.encode([question], convert_to_numpy=True)

#     # --- FAISS + BM25 retrieval
#     faiss_idx = index.search(q_emb, 10)[1][0]
#     bm25_idx = np.argsort(BM25.get_scores(re.findall(r"[a-z0-9]+", question.lower())))[::-1][:10]

#     # --- Reciprocal Rank Fusion (RRF)
#     rrf_scores = {}
#     for rank, i in enumerate(faiss_idx): rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
#     for rank, i in enumerate(bm25_idx): rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
#     fused = sorted(rrf_scores, key=rrf_scores.get, reverse=True)[:10]

#     # --- Cross-encoder rerank
#     pairs = [(question, METADATA[i]["text"]) for i in fused]
#     scores = RERANKER.predict(pairs)
#     ranked = [fused[i] for i in np.argsort(-scores)[:5]]

#     # --- LLM synthesis
#     context = "\n---\n".join(METADATA[i]["text"] for i in ranked)
#     model = genai.GenerativeModel(GEMINI_MODEL)
#     prompt = f"Role: {role}\nUse the context below to answer clearly.\n\nCONTEXT:\n{context}\n\nQUESTION:\n{question}\n\nANSWER:"
#     try:
#         print("\n==============================")
#         print("📤 Gemini Request Sent")
#         print("==============================")
#         print(prompt)
#         print("==============================")

#         response = model.generate_content(prompt)
#         print("\n✅ Gemini Response Received")
#         print("==============================")
#         print(response.text)
#         print("==============================\n")
        
        
#         return response.text.strip()
#     except Exception as e:
#         return f"⚠️ LLM generation failed: {e}"




# # def answer_rerank_only(question: str):
# #     """Retrieval-only mode: FAISS + BM25 + RRF + Reranker"""
# #     q_emb = EMB.encode([question], convert_to_numpy=True)

# #     faiss_idx = index.search(q_emb, 10)[1][0]
# #     bm25_idx = np.argsort(BM25.get_scores(re.findall(r"[a-z0-9]+", question.lower())))[::-1][:10]

# #     rrf_scores = {}
# #     for rank, i in enumerate(faiss_idx): rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
# #     for rank, i in enumerate(bm25_idx): rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
# #     fused = sorted(rrf_scores, key=rrf_scores.get, reverse=True)[:10]

# #     pairs = [(question, METADATA[i]["text"]) for i in fused]
# #     scores = RERANKER.predict(pairs)
# #     ranked = [fused[i] for i in np.argsort(-scores)[:3]]

# #     chunks = [METADATA[i]["text"] for i in ranked]
# #     answer = " ".join(ch[:200] for ch in chunks)
# #     return answer.strip()

# def answer_rerank_only(question: str, max_sentences: int = 3):
#     """
#     Pure retrieval-only mode (no LLM)
#     Returns: Short direct answer + inline citations like [1][2]
#     """
#     import re
#     import numpy as np

#     if len(METADATA) == 0:
#         return "No documents uploaded yet."

#     # Step 1: Hybrid retrieval (FAISS + BM25 + RRF)
#     q_emb = EMB.encode([question], convert_to_numpy=True)
#     faiss_idx = index.search(q_emb, 20)[1][0]
#     bm25_scores = BM25.get_scores(re.findall(r"[a-z0-9]+", question.lower()))
#     bm25_idx = np.argsort(bm25_scores)[::-1][:20]

#     rrf_scores = {}
#     for rank, i in enumerate(faiss_idx):
#         if i >= 0: rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
#     for rank, i in enumerate(bm25_idx):
#         rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)

#     fused_chunk_ids = sorted(rrf_scores.items(), key=lambda x: x[1], reverse=True)[:12]
#     chunk_ids = [cid for cid, _ in fused_chunk_ids]

#     # Step 2: Split chunks into sentences + track source
#     sentences = []
#     citations = []  # Will store [source, page] for each sentence

#     for chunk_id in chunk_ids:
#         text = METADATA[chunk_id]["text"]
#         source = METADATA[chunk_id]["source"]
#         page = METADATA[chunk_id]["page"]

#         # Smart sentence split
#         raw_sents = re.split(r'(?<=[.!?])\s+', text.strip())
#         for sent in raw_sents:
#             sent = sent.strip()
#             if len(sent) > 20 and sent not in ["•", "-", "–", "—"]:
#                 sentences.append(sent)
#                 citations.append({
#                     "source": source,
#                     "page": page or "-",
#                     "chunk_id": chunk_id
#                 })

#     if not sentences:
#         return "I couldn't find a clear answer in the documents."

#     # Step 3: Rerank sentences with cross-encoder
#     pairs = [(question, sent) for sent in sentences]
#     scores = RERANKER.predict(pairs)
#     top_indices = np.argsort(scores)[-max_sentences:][::-1]

#     # Step 4: Build answer with inline citations
#     answer_parts = []
#     used_refs = set()

#     for idx in top_indices:
#         sent = sentences[idx]
#         ref = citations[idx]
#         ref_key = (ref["source"], ref["page"])

#         if ref_key not in used_refs:
#             used_refs.add(ref_key)
#             ref_num = len(used_refs)
#             answer_parts.append(f"{sent}[{ref_num}]")
#         else:
#             ref_num = list(used_refs).index(ref_key) + 1
#             answer_parts.append(f"{sent}[{ref_num}]")

#     answer = " ".join(answer_parts)

#     # Step 5: Add citation list at the end
#     citation_list = ""
#     for i, (src, pg) in enumerate(used_refs, 1):
#         citation_list += f"\n[{i}] {Path(src).name}" + (f" (p.{pg})" if pg != "-" else "")

#     final_answer = answer.strip()
#     if citation_list:
#         final_answer += "\n\nSources:" + citation_list

#     return final_answer


# def answer_simple_retrieval(question: str):
#     """Simple mode: FAISS + BM25 + RRF (no reranker, no LLM)"""
#     q_emb = EMB.encode([question], convert_to_numpy=True)

#     faiss_idx = index.search(q_emb, 10)[1][0]
#     bm25_idx = np.argsort(BM25.get_scores(re.findall(r"[a-z0-9]+", question.lower())))[::-1][:10]

#     rrf_scores = {}
#     for rank, i in enumerate(faiss_idx): rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
#     for rank, i in enumerate(bm25_idx): rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
#     fused = sorted(rrf_scores, key=rrf_scores.get, reverse=True)[:5]

#     chunks = [METADATA[i]["text"] for i in fused]
#     answer = " ".join(ch[:200] for ch in chunks)
#     return answer.strip()




# ==================== THREE MODES FOR RETRIEVAL (Replace your answer functions with these)
# def answer_hybrid_llm(question: str, role: str = "User"):
#     """Mode 1: Hybrid LLM — Question-related answer with context awareness"""
#     q_emb = EMB.encode([question], convert_to_numpy=True)

#     # Retrieval
#     faiss_idx = index.search(q_emb, 12)[1][0]
#     bm25_idx = np.argsort(BM25.get_scores(re.findall(r"[a-z0-9]+", question.lower())))[::-1][:12]

#     rrf_scores = {}
#     for rank, i in enumerate(faiss_idx): rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
#     for rank, i in enumerate(bm25_idx): rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
#     fused = sorted(rrf_scores, key=rrf_scores.get, reverse=True)[:10]

#     # Rerank
#     pairs = [(question, METADATA[i]["text"]) for i in fused]
#     scores = RERANKER.predict(pairs)
#     ranked = [fused[i] for i in np.argsort(-scores)[:6]]

#     # LLM for precise synthesis
#     context = "\n---\n".join(METADATA[i]["text"] for i in ranked)
#     model = genai.GenerativeModel(GEMINI_MODEL)
#     prompt = f"Role: {role}\nFocus ONLY on question-related info from context. Be concise.\n\nCONTEXT:\n{context}\n\nQUESTION:\n{question}\n\nANSWER:"
#     try:
#         response = model.generate_content(prompt)
#         return response.text.strip()
#     except Exception as e:
#         return f"Error: {e}"


from groq import Groq
import os
import re
import numpy as np

# Initialize Groq client (set your API key in environment or directly)
client = Groq(api_key=os.getenv("GROQ_API_KEY"))  # Recommended: use env variable

# Make sure these are defined elsewhere in your code:
# EMB, index, BM25, METADATA, RERANKER

def answer_hybrid_llm(question: str, role: str = "User"):
    """Mode 1: Hybrid LLM — Now powered by Groq Llama-3.1-8B-Instant (super fast!)"""
    q_emb = EMB.encode([question], convert_to_numpy=True)

    # Retrieval: FAISS + BM25
    faiss_idx = index.search(q_emb, 12)[1][0]
    bm25_idx = np.argsort(BM25.get_scores(re.findall(r"[a-z0-9]+", question.lower())))[::-1][:12]

    # Reciprocal Rank Fusion (RRF)
    rrf_scores = {}
    for rank, i in enumerate(faiss_idx):
        rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
    for rank, i in enumerate(bm25_idx):
        rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
    
    fused = sorted(rrf_scores, key=rrf_scores.get, reverse=True)[:10]

    # Reranking (using cross-encoder)
    pairs = [(question, METADATA[i]["text"]) for i in fused]
    scores = RERANKER.predict(pairs)
    ranked = [fused[i] for i in np.argsort(-scores)[:4]]

    # Build context
    context = "\n---\n".join(METADATA[i]["text"] for i in ranked)
    # Prompt for Llama-3.1 via Groq
    prompt = f"""You are a helpful and precise assistant. Answer the question using ONLY the provided context. 
If the context doesn't contain enough information, say "I don't have sufficient information to answer accurately."

Role: {role}

Context:
{context}


Question: {question}

Answer:"""

    try:
        chat_completion = client.chat.completions.create(
            messages=[
                {"role": "system", "content": "You are a knowledgeable assistant. Be concise, accurate, and professional."},
                {"role": "user", "content": prompt}
            ],
            model="llama-3.1-8b-instant",   # Super fast & cheap
            # model="llama-3.1-70b-versatile",  # Use this if you want better reasoning (more expensive)
            temperature=0.3,
            max_tokens=512,
            top_p=0.95,
        )
       
        print("📝 Groq Llama Response:\n", chat_completion.choices[0].message.content.strip())  # <--- DEBUG RESPONSE
        response = chat_completion.choices[0].message.content.strip()
        return response

    except Exception as e:
        return f"Groq API Error: {str(e)}"




def answer_rerank_only(question: str, max_sentences: int = 3):
    """Mode 2: Pure Retrieval — LLM-level accuracy without LLM"""
    import numpy as np

    q_emb = EMB.encode([question], convert_to_numpy=True)
    faiss_idx = index.search(q_emb, 15)[1][0]
    bm25_idx = np.argsort(BM25.get_scores(re.findall(r"[a-z0-9]+", question.lower())))[::-1][:15]

    rrf_scores = {}
    for rank, i in enumerate(faiss_idx): rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
    for rank, i in enumerate(bm25_idx): rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
    fused = sorted(rrf_scores, key=rrf_scores.get, reverse=True)[:10]

    # Sentence-level extraction
    sentences = []
    for i in fused:
        text = METADATA[i]["text"]
        sents = re.split(r'(?<=[.!?])\s+', text.strip())
        sentences.extend([s.strip() for s in sents if len(s) > 20])

    if not sentences:
        return "No relevant answer found."

    pairs = [(question, s) for s in sentences]
    scores = RERANKER.predict(pairs)
    top_indices = np.argsort(scores)[-max_sentences:][::-1]
    best_sentences = [sentences[j] for j in top_indices]

    answer = " ".join(best_sentences)
    return answer.strip()

def answer_simple_retrieval(question: str):
    """Mode 3: Simple Retrieval — Basic question-related answer"""
    q_emb = EMB.encode([question], convert_to_numpy=True)
    faiss_idx = index.search(q_emb, 10)[1][0]
    bm25_idx = np.argsort(BM25.get_scores(re.findall(r"[a-z0-9]+", question.lower())))[::-1][:10]

    rrf_scores = {}
    for rank, i in enumerate(faiss_idx): rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
    for rank, i in enumerate(bm25_idx): rrf_scores[i] = rrf_scores.get(i, 0) + 1 / (60 + rank)
    fused = sorted(rrf_scores, key=rrf_scores.get, reverse=True)[:5]

    chunks = [METADATA[i]["text"] for i in fused]
    answer = " ".join(ch[:200] for ch in chunks if ch.lower().find(question.lower()) > -1)  # Filter question-related
    return answer.strip() or "No direct match found."






# ----------------------------------------------------------
# ✅ FastAPI Server
# ----------------------------------------------------------
app = FastAPI(title="Gemini RAG Backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/debug")
def debug():
    """Shows whether Gemini is available & model ID."""
    return {
        "gemini_key_loaded": bool(GEMINI_API_KEY),
        "gemini_model": GEMINI_MODEL,
        "chunks_indexed": len(METADATA),
        "files_count": len(set([m["source"] for m in METADATA])),
    }




# ===============================================
# 📊 Analytics Dashboard Metrics (backend)
# ===============================================
# from fastapi import Request
# import time
# from statistics import mean

# # Store analytics data in memory (or switch to DB later)
# ANALYTICS = {
#     "search_count": 0,
#     "sessions": set(),
#     "search_times": [],
#     "embed_times": [],
#     "queries": {}
# }




# @app.get("/metrics")
# async def get_metrics():
#     if METRICS["search_count"] == 0:
#         return {
#             "search_count": 0,
#             "avg_time_ms": 0,
#             "avg_embed_ms": 0,
#             "avg_llm_ms": 0,
#             "sessions": 0,
#             "top_queries": [],
#             "last_ingest_at": None,
#             "last_persist_at": None,
#         }

#     top = sorted(METRICS["queries"].items(), key=lambda x: x[1], reverse=True)[:5]

#     return {
#         "search_count": METRICS["search_count"],
#         "avg_time_ms": round(mean(METRICS["total_ms"]), 2),
#         "avg_embed_ms": round(mean(METRICS["embed_ms"]), 2),
#         "avg_llm_ms": round(mean(METRICS["llm_ms"]), 2),
#         "sessions": len(METRICS["sessions"]),
#         "top_queries": top,
#         "last_ingest_at": METRICS["last_updated"],
#         "last_persist_at": METRICS["last_updated"],
#     }
















OCR_USED = False   # ← will become True once any image is processed

@app.get("/health")
def health():
    return {
        "status": "🟢 online",
        "chunks": len(METADATA),
        "files": len(set(m["source"] for m in METADATA)),
        "emb_model": index.ntotal,
        "bm25": len(token_corpus),
        "ocr": OCR_USED,
        # Add this line → frontend will use it
        "display": (
            f"Status: online\n"
            f"Chunks: {len(METADATA)}\n"
            f"Files: {len(set(m['source'] for m in METADATA))}\n"
            f"Embeddings: {index.ntotal}\n"
            f"BM25 docs: {len(token_corpus)}\n"
            # f"OCR: {'Yes' if False else 'No'}"
            f"OCR: {'Enabled (used)' if OCR_USED else 'Enabled (not used yet)'}"
        )
    }




@app.get("/chunks")
def chunks():
    return [
        {
            "id": i,
            "page": m["page"],
            "source": m["source"],
            "text": m["text"]
        }
        for i, m in enumerate(METADATA)
    ]









@app.post("/reset")
def reset():
    shutil.rmtree(DATA, ignore_errors=True)
    DATA.mkdir(exist_ok=True)

    global index, METADATA, token_corpus, BM25
    index = faiss.IndexFlatL2(DIM)
    METADATA = []
    token_corpus = []
    BM25 = None
    persist()

    return {"message": "Reset OK"}


@app.get("/files")
def files():
    return {"files": sorted(set([m["source"] for m in METADATA]))}




@app.post("/upload")
async def upload(files: List[UploadFile] = File(...)):
    global OCR_USED
    total = 0

    for f in files:

        # ✅ read file bytes
        raw = await f.read()

        # ✅ save uploaded file
        UPLOADS.joinpath(f.filename).write_bytes(raw)

        # ✅ extract text blocks directly from raw bytes (NO UploadFile recreation)
        name = f.filename.lower()
        if name.endswith(".pdf"):
            blocks = extract_pdf(raw, f.filename)
        elif name.endswith(".docx"):
            blocks = extract_docx(raw, f.filename)
        elif name.endswith(".txt"):
            blocks = extract_txt(raw, f.filename)
        elif name.endswith((".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp")):
            OCR_USED = True   # ← THIS LINE MAKES IT DYNAMIC
            blocks = extract_image(raw, f.filename)   # OCR kicks in here!
        else:
            blocks = extract_txt(raw, f.filename)

        # ✅ ingest into FAISS + BM25
        total += ingest(blocks)

    return {
        "status": "ok",
        "chunks": len(METADATA),
        "files": len(set([m["source"] for m in METADATA])),
        "emb_model": index.ntotal,
        "bm25": len(token_corpus),
        "ocr": False
    }



import time
from statistics import mean
from datetime import datetime

METRICS = {
    "search_count": 0,
    "total_ms": [],
    "embed_ms": [],
    "llm_ms": [],
    "sessions": set(),
    "queries": {},
    "last_updated": None
}

@app.post("/chat")
async def chat(request: Request, payload: dict):
    import time
    user_id = get_user_id(request)
    if user_id and hasattr(request.state, "user") and request.state.user:
        user_id = request.state.user.get("email") or user_id
    if not user_id:
        return {"reply": "Unauthorized", "mode": "error"}

    question = payload["messages"][-1]["content"]
    chat_id = payload.get("chat_id") or str(uuid.uuid4())
    role = payload.get("role", "User")
    mode = payload.get("retrieval_mode", "hybrid")

    total_start = time.time()
    embed_start = time.time()

    # ——— RETRIEVAL + ANSWER ———
    if mode == "hybrid":
        # Get ranked chunk IDs
        ranked_ids = rerank(question, topk=3)  # ← uses your existing rerank() function

        # Build context for LLM
        context_chunks = [METADATA[i] for i in ranked_ids]
        context = "\n---\n".join(chunk["text"] for chunk in context_chunks)

        # Call Groq Llama
        reply = answer_hybrid_llm(question, role)

        # ——— BUILD DETAILED SOURCES FOR FRONTEND ———
        sources = []
        for rank, chunk in enumerate(context_chunks, 1):
            sources.append({
                "rank": rank,
                "text": chunk["text"].strip(),
                "source": chunk["source"],
                "page": chunk["page"],
                "type": "OCR" if str(chunk["source"]).lower().endswith(('.png', '.jpg', '.jpeg')) else "Text"
            })

    else:
        reply = "Unknown mode"
        sources = []

    # ←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←
    # CALCULATE LATENCIES (this was completely broken before)
    # ——— ACCURATE LATENCY MEASUREMENT ———
    embed_start = time.time()
    ranked_ids = rerank(question, topk=5)                    # ← embedding + rerank happens here
    embed_ms = round((time.time() - embed_start) * 1000, 1)

    llm_start = time.time()
    reply = answer_hybrid_llm(question, role)                # ← only LLM happens here
    llm_ms = round((time.time() - llm_start) * 1000, 1)

    total_ms = round(embed_ms + llm_ms, 1)                   # real total

    # ← accurate total       # rough but works perfectly
    # ←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←←

    # Update global metrics
    METRICS["search_count"] += 1
    METRICS["total_ms"].append(total_ms)
    METRICS["embed_ms"].append(embed_ms)
    METRICS["llm_ms"].append(llm_ms)
    METRICS["sessions"].add(request.headers.get("X-Session-ID", "unknown"))
    METRICS["queries"][question] = METRICS["queries"].get(question, 0) + 1
    METRICS["last_updated"] = datetime.utcnow().isoformat()   # keep this
    METRICS["last_ingest_at"] = METRICS["last_updated"]       # ← ADD THIS LINE
    METRICS["last_persist_at"] = METRICS["last_updated"]      # ← AND THIS    # Save chat to Redis
    if r:
        key = f"chat:{user_id}:{chat_id}"
        r.rpush(key, json.dumps({"role": "user", "content": question}))
        r.rpush(key, json.dumps({"role": "assistant", "content": reply, "sources": sources}))
        r.expire(key, 60 * 60 * 24 * 90)

        title_key = f"chat_title:{user_id}:{chat_id}"
        if not r.exists(title_key):
            short_title = question.strip()[:40] + ("..." if len(question) > 40 else "")
            r.set(title_key, short_title)

    return {
        "reply": reply,
        "mode": mode,
        "chat_id": chat_id,
        "sources": sources  # ← THIS IS WHAT YOUR FRONTEND NEEDS!
    }

@app.get("/chats")
async def list_chats(request: Request):  # Note: /chats (plural) for list
    user_id = get_user_id(request)
    if not user_id or not r: 
        return {"chats": []}
    keys = r.keys(f"chat:{user_id}:*")
    chats = []
    for k in keys:
        chat_id = k.split(":")[-1]  # ← Remove .decode() — k is already str
        messages = r.lrange(k, 0, 0)
        title = "New Chat"
        if messages:
            first = json.loads(messages[0])  # ← Remove .decode() — messages[0] is already str
            title = first["content"][:30] + "..." if len(first["content"]) > 30 else first["content"]
        chats.append({"id": chat_id, "title": title})
    return {"chats": sorted(chats, key=lambda x: x["id"], reverse=True)}


@app.post("/chats/new")
async def create_new_chat(request: Request):
    # ← THIS ONE LINE FIXES EVERYTHING (uses same logic as your working get_user_id)
    user_id = get_user_id(request)
    if user_id and hasattr(request.state, "user") and request.state.user:
        user_id = request.state.user.get("email") or user_id
        
    if not user_id:
        return JSONResponse({"error": "Unauthorized"}, status_code=401)

    # ← NOW FORCE IT TO EMAIL LIKE YOU DID IN /chat POST
    if hasattr(request.state, "user") and request.state.user:
        user_id = request.state.user.get("email") or user_id

    chat_id = str(uuid.uuid4())
    key = f"chat:{user_id}:{chat_id}"
    
    # Create empty chat so it exists
    r.rpush(key, json.dumps({"role": "system", "content": "Chat started"}))
    r.expire(key, 60 * 60 * 24 * 90)

    return {"chat_id": chat_id, "title": "New Chat"}


@app.get("/chats/{chat_id}")
async def get_chat(chat_id: str, request: Request):
    # ← 100% SAFE user_id extraction – never triggers 422
    auth_header = request.headers.get("Authorization")
    user_id = None
    if auth_header and auth_header.startswith("Bearer "):
        try:
            token = auth_header.split(" ")[1]
            payload = jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
            user_id = payload.get("sub")  # this is the email
        except:
            pass

    if not user_id:
        return JSONResponse({"error": "Unauthorized"}, status_code=401)

    key = f"chat:{user_id}:{chat_id}"
    if not r or not r.exists(key):
        return JSONResponse({"error": "Chat not found"}, status_code=404)

    raw = r.lrange(key, 0, -1)
    messages = [json.loads(m) for m in raw]

    title = "New Chat"
    title_key = f"chat_title:{user_id}:{chat_id}"
    if r.exists(title_key):
        title = r.get(title_key) or "New Chat"

    response = {
        "id": chat_id,
        "title": title,
        "messages": messages
    }
    print("RESPONSE →", json.dumps(response, indent=2))
    return response

@app.delete("/chats/clear-all")
async def clear_all_chats(request: Request):
    user_id = get_user_id(request)
    if not user_id or not r:
        return {"ok": False}
    keys = r.keys(f"chat:{user_id}:*")
    if keys:
        r.delete(*keys)
    return {"ok": True}


# In-memory metrics (add near the top with other globals)
METRICS = {
    "search_count": 0,
    "total_ms": [], "embed_ms": [], "llm_ms": [],
    "sessions": set(),
    "queries": {},
    "last_updated": None
}



# DELETE BOTH OLD /metrics ENDPOINTS
# PUT THIS ONE AT THE VERY BOTTOM OF THE FILE (after @app.get("/me"))

@app.get("/metrics")
async def get_metrics():
    if METRICS["search_count"] == 0:
        return {
            "search_count": 0,
            "avg_time_ms": 0,
            "avg_embed_ms": 0,
            "avg_llm_ms": 0,
            "sessions": 0,
            "top_queries": [],
            "last_ingest_at": None,
            "last_persist_at": None,
        }

    from statistics import mean
    top = sorted(METRICS["queries"].items(), key=lambda x: x[1], reverse=True)[:8]

    return {
        "search_count": METRICS["search_count"],
        "avg_time_ms": round(mean(METRICS["total_ms"]), 1),
        "avg_embed_ms": round(mean(METRICS["embed_ms"]), 1),
        "avg_llm_ms": round(mean(METRICS["llm_ms"]), 1),
        "sessions": len(METRICS["sessions"]),
        "top_queries": top,
        "last_ingest_at": METRICS["last_updated"],
        "last_persist_at": METRICS["last_updated"],
    }


from starlette.middleware.sessions import SessionMiddleware
from fastapi.responses import RedirectResponse, JSONResponse
from jose import jwt, JWTError
from datetime import datetime, timedelta

# Add session support (needed for OAuth flow)
app.add_middleware(SessionMiddleware, secret_key="your-super-secret-key-change-in-production")

# Update CORS to allow your frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5500", "http://127.0.0.1:5500"],  # Both work
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Google OAuth
GOOGLE_CLIENT_ID = os.getenv("GOOGLE_CLIENT_ID")
GOOGLE_CLIENT_SECRET = os.getenv("GOOGLE_CLIENT_SECRET")

if not GOOGLE_CLIENT_ID or not GOOGLE_CLIENT_SECRET:
    print("⚠️ WARNING: GOOGLE_CLIENT_ID or GOOGLE_CLIENT_SECRET missing in .env")



oauth = OAuth()

# ========================================
# GOOGLE OAUTH – FINAL WORKING VERSION
# ========================================
from authlib.integrations.starlette_client import OAuth

oauth = OAuth()

# THIS IS THE ONLY ONE THAT MUST EXIST
oauth.register(
    name='google',
    client_id=os.getenv("GOOGLE_CLIENT_ID"),
    client_secret=os.getenv("GOOGLE_CLIENT_SECRET"),
    server_metadata_url='https://accounts.google.com/.well-known/openid-configuration',
    client_kwargs={
        'scope': 'openid email profile'
    }
)

# Force consent screen → guarantees id_token every time
@app.get("/auth/google/login")
async def google_login(request: Request):
    redirect_uri = "http://localhost:8000/auth/google/callback"
    return await oauth.google.authorize_redirect(
        request,
        redirect_uri,
        prompt="consent"      # ← this now actually works!
    )


@app.get("/auth/google/callback")
async def google_callback(request: Request):
    try:
        token = await oauth.google.authorize_access_token(request)

        # Method 1: Try to get id_token directly
        user = token.get("userinfo")
        if not user:
            # Method 2: Try parse_id_token (might still fail)
            try:
                user = await oauth.google.parse_id_token(request, token)
            except:
                user = None

        # Method 3: FINAL FALLBACK — fetch userinfo with access_token
        if not user and "access_token" in token:
            async with oauth.google.create_client_session() as session:
                resp = await session.get(
                    "https://www.googleapis.com/oauth2/v3/userinfo",
                    headers={"Authorization": f"Bearer {token['access_token']}"}
                )
                resp.raise_for_status()
                user = resp.json()

        if not user or "email" not in user:
            return JSONResponse({"error": "Failed to get user info from Google"}, status_code=400)

        # Success! Create your JWT
        jwt_token = create_jwt_token({
            "email": user["email"],
            "name": user.get("name"),
            "picture": user.get("picture"),
        })

        return RedirectResponse(f"http://localhost:5500/ui/index.html?jwt={jwt_token}&name={user.get('name','')}")

    except Exception as e:
        print("Google OAuth Error:", e)
        import traceback
        traceback.print_exc()
        return JSONResponse({"error": "Authentication failed"}, status_code=400)

# JWT Config
JWT_SECRET = "your-jwt-secret-key-please-change-this-too"
JWT_ALGORITHM = "HS256"
JWT_EXPIRE_MINUTES = 60

def create_jwt_token(user_data: dict):
    expire = datetime.utcnow() + timedelta(minutes=JWT_EXPIRE_MINUTES)
    payload = {
        "sub": user_data["email"],
        "name": user_data.get("name"),
        "picture": user_data.get("picture"),
        "exp": expire
    }
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGORITHM)


@app.get("/me")
async def get_current_user(Authorization: str = None):
    if not Authorization:
        raise HTTPException(status_code=401, detail="No token")
    try:
        scheme, _, token = Authorization.partition(" ")
        payload = jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        return payload
    except JWTError:
        raise HTTPException(status_code=401, detail="Invalid token")


if __name__ == "__main__":
    print("✅ Gemini RAG running: http://localhost:8000")
    uvicorn.run("new:app", host="0.0.0.0", port=8000, reload=True)
